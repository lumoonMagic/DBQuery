import os
from datetime import datetime
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class BoardExporter:
    def __init__(self, demo_mode=True):
        self.demo_mode = demo_mode
        self.pr = Presentation()
        self.slide_layout = self.pr.slide_layouts[5]  # blank slide

    def add_title_slide(self, title="Board Report", subtitle=None):
        slide = self.pr.slides.add_slide(self.slide_layout)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = title
        p.font.size = Pt(36)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        if subtitle:
            p2 = text_frame.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.font.italic = True
            p2.alignment = PP_ALIGN.CENTER

    def add_table_slide(self, title, dataframe):
        slide = self.pr.slides.add_slide(self.slide_layout)
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(0.5))
        textbox.text = title

        rows, cols = dataframe.shape
        left = Inches(0.5)
        top = Inches(1)
        width = Inches(9)
        height = Inches(0.8 + 0.2 * rows)

        table_shape = slide.shapes.add_table(rows+1, cols, left, top, width, height)
        table = table_shape.table

        # Set header
        for j, col_name in enumerate(dataframe.columns):
            table.cell(0, j).text = str(col_name)
            table.cell(0, j).text_frame.paragraphs[0].font.bold = True
            table.cell(0, j).text_frame.paragraphs[0].font.size = Pt(12)

        # Set data
        for i in range(rows):
            for j in range(cols):
                table.cell(i+1, j).text = str(dataframe.iloc[i, j])
                table.cell(i+1, j).text_frame.paragraphs[0].font.size = Pt(10)

    def add_text_slide(self, title, text):
        slide = self.pr.slides.add_slide(self.slide_layout)
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(6))
        tf = textbox.text_frame
        tf.text = f"{title}\n{text}"
        for p in tf.paragraphs:
            p.font.size = Pt(14)

    def save_ppt(self, output_path=None):
        if self.demo_mode:
            output_path = output_path or f"demo_board_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        else:
            output_path = output_path or f"board_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

        self.pr.save(output_path)
        return output_path

# --------------------------
# Example usage
# --------------------------
if __name__ == "__main__":
    import pandas as pd
    exporter = BoardExporter(demo_mode=True)
    exporter.add_title_slide("Demo Vendor Analysis", "Supply Chain Insights")
    df = pd.DataFrame({
        "Vendor": ["V1", "V2", "V3"],
        "OTIF": [98, 92, 85],
        "Defect Rate": [0.5, 1.2, 2.3]
    })
    exporter.add_table_slide("Vendor Performance Table", df)
    exporter.add_text_slide("Suggested Improvements", "Focus on vendor V3 for on-time improvement initiatives.")
    output_file = exporter.save_ppt()
    print(f"Demo PPT saved at: {output_file}")
